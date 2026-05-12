"""
EL 장애분석 보고서 생성 (python-pptx 기반)
Vercel Serverless 환경 - 외부 프로세스 없음
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io

def rgb(hex_str):
    h=hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# 색상
NAVY  = "0F2149"; BLUE  = "1D4ED8"; SKY   = "3B82F6"
SKYLT = "93C5FD"; ICE   = "DBEAFE"; TEAL  = "0D7966"
TEALLT= "14B8A6"; AMBER = "F59E0B"; GREEN = "059669"
RED   = "DC2626"; DARK  = "1E293B"; GRAY  = "64748B"
LGRAY = "E2E8F0"; WHITE = "FFFFFF"; OFFWH = "F8FAFC"

W = Inches(10); H = Inches(5.625)

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        if line_w: shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Malgun Gothic", wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.color.rgb = rgb(color)
    run.font.bold = bold; run.font.italic = italic
    run.font.name = font
    return txBox

def add_header(slide, num, title, subtitle=None):
    """기존 PPT 스타일 헤더"""
    add_rect(slide, 0, 0, 10, 0.68, NAVY)
    add_rect(slide, 0, 0, 0.07, 0.68, AMBER)
    add_text(slide, num, 0.15, 0.05, 0.55, 0.58, 10, AMBER, True, PP_ALIGN.CENTER, "Calibri")
    add_text(slide, title, 0.75, 0.08, 8.5, 0.55, 19, WHITE, True)
    if subtitle:
        add_text(slide, subtitle, 6.5, 0.08, 3.4, 0.55, 10, SKYLT, False, PP_ALIGN.RIGHT)
    add_rect(slide, 0, 0.68, 10, 0.045, BLUE)

def card(slide, x, y, w, h, accent=None):
    add_rect(slide, x, y, w, h, WHITE, LGRAY, 0.5)
    if accent:
        add_rect(slide, x, y, w, 0.06, accent)

def build_el_report(stats):
    """EL 장애분석 보고서 생성 → bytes 반환"""
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    s1el = stats.get('1호선EL', 0)
    s2el = stats.get('2호선EL', 0)
    total_el = s1el + s2el

    # ── 슬라이드 1: 표지 ──────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, 10, 5.625, NAVY)
    add_rect(sl, 0, 0, 10, 0.07, AMBER)
    add_rect(sl, 0, 0.07, 0.5, 5.555, BLUE)
    # 원형 장식
    for ox,oy,ow,oh,oc,tr in [(6.2,-0.3,4.5,4.5,BLUE,85),(7.0,0.5,3.0,3.0,SKY,80),(7.5,1.0,2.0,2.0,SKYLT,70)]:
        sh = sl.shapes.add_shape(9, Inches(ox), Inches(oy), Inches(ow), Inches(oh))
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(oc)
        from pptx.util import Pt as P2
        sh.line.color.rgb = rgb(WHITE); sh.line.width = P2(0.5)
    add_rect(sl, 0.7, 1.2, 1.6, 0.38, BLUE)
    add_text(sl, "ELEVATOR  ·  EL", 0.7, 1.22, 1.6, 0.34, 9, WHITE, True, PP_ALIGN.CENTER, "Calibri")
    add_rect(sl, 0.7, 1.65, 2.4, 0.34, AMBER)
    add_text(sl, "2025년  4분기", 0.7, 1.67, 2.4, 0.3, 12, NAVY, True, PP_ALIGN.CENTER)
    add_text(sl, "엘리베이터", 0.7, 2.12, 9, 0.95, 46, WHITE, True, PP_ALIGN.LEFT, "Calibri")
    add_text(sl, "장애분석 보고서", 0.7, 3.0, 9, 0.9, 38, SKYLT, True, PP_ALIGN.LEFT, "Calibri")
    add_text(sl, "1호선 · 2호선  |  대구도시철도공사 시설운영처", 0.7, 4.0, 8.5, 0.38, 12, "8EB8D8")
    add_rect(sl, 0, 5.22, 10, 0.405, "070F2B")
    add_text(sl, f"EL 1호선 {s1el}건 · EL 2호선 {s2el}건  |  합계 {total_el}건",
             0.5, 5.24, 9, 0.35, 10.5, "6B9AC4", False, PP_ALIGN.CENTER)

    # ── 슬라이드 2: 종합 현황 ─────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl, "01", "종합 현황", "2025년 4분기 (10월 ~ 12월)")
    add_rect(sl, 0, 0.69, 10, 0.04, BLUE)

    kpis = [
        ("1호선 EL 장애건수", str(s1el)+"건", "전분기 37건 대비 +22건", BLUE, RED),
        ("2호선 EL 장애건수", str(s2el)+"건", "전분기 27건 대비 +7건", BLUE, RED),
        ("평균 복구시간", "91분", "1호선 102분 · 2호선 74분", GRAY, GRAY),
        ("총 미가동시간", "116h", "1호선 77h · 2호선 38h", AMBER, GREEN),
    ]
    for i,(label,val,sub,c,tc) in enumerate(kpis):
        x=0.2+i*2.45
        card(sl, x, 0.85, 2.35, 1.52, c)
        add_text(sl,label,x,0.93,2.35,0.28,9.5,c,True,PP_ALIGN.CENTER)
        add_text(sl,val,x,1.18,2.35,0.65,30,DARK,True,PP_ALIGN.CENTER,"Calibri")
        add_text(sl,sub,x,1.83,2.35,0.28,8.5,tc,False,PP_ALIGN.CENTER)

    # 1호선 상세
    card(sl, 0.2, 2.5, 4.65, 2.82, BLUE)
    add_text(sl,"1호선 EL  —  "+str(s1el)+"건",0.28,2.58,4.5,0.32,12,BLUE,True)
    ranks=[("안전스위치","14건","23.7%"),("전장품","13건","22.0%"),("도어","11건","18.6%")]
    for j,(n,v,p) in enumerate(ranks):
        ry=2.93+j*0.66
        add_rect(sl,2.5,ry,2.25,0.58,ICE if j==0 else WHITE,LGRAY,0.3)
        bc=BLUE if j==0 else LGRAY
        sh=sl.shapes.add_shape(9,Inches(2.55),Inches(ry+0.1),Inches(0.35),Inches(0.35))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(bc); sh.line.fill.background()
        add_text(sl,str(j+1),2.55,ry+0.1,0.35,0.35,10,WHITE,True,PP_ALIGN.CENTER)
        add_text(sl,n,2.95,ry+0.02,1.75,0.28,10,DARK,j==0)
        add_text(sl,v+"("+p+")",2.95,ry+0.3,1.75,0.22,8.5,GRAY)

    # 2호선 상세
    card(sl, 5.1, 2.5, 4.7, 2.82, SKY)
    add_text(sl,"2호선 EL  —  "+str(s2el)+"건",5.18,2.58,4.55,0.32,12,BLUE,True)
    ranks2=[("안전스위치","21건","61.8%"),("도어","5건","14.7%"),("속도조정장치","3건","8.8%")]
    for j,(n,v,p) in enumerate(ranks2):
        ry=2.93+j*0.66
        add_rect(sl,7.38,ry,2.28,0.58,ICE if j==0 else WHITE,LGRAY,0.3)
        bc=SKY if j==0 else LGRAY
        sh=sl.shapes.add_shape(9,Inches(7.43),Inches(ry+0.1),Inches(0.35),Inches(0.35))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(bc); sh.line.fill.background()
        add_text(sl,str(j+1),7.43,ry+0.1,0.35,0.35,10,WHITE,True,PP_ALIGN.CENTER)
        add_text(sl,n,7.83,ry+0.02,1.78,0.28,10,DARK,j==0)
        add_text(sl,v+"("+p+")",7.83,ry+0.3,1.78,0.22,8.5,GRAY)

    add_rect(sl,0.2,5.35,9.6,0.22,"FEF3C7","F59E0B",0.5)
    add_text(sl,"⚠  2호선 EL 안전스위치 61.8%  |  1호선 EL 안전스위치·전장품·도어 3대 장치가 64.3% 차지",
             0.25,5.35,9.5,0.22,9,"92400E",True,PP_ALIGN.LEFT)

    # ── 슬라이드 3: 월별 현황 ─────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl, "02", "월별 장애 현황", "10월 · 11월 · 12월")
    card(sl, 0.2, 0.85, 9.6, 3.2, BLUE)
    add_text(sl,"EL 1·2호선 월별 장애건수 추이",0.28,0.93,6,0.32,12,BLUE,True)
    add_text(sl,"1호선: 13→21→25건 (증가)  ·  2호선: 12→9→13건",6.0,0.93,3.75,0.32,9,GRAY,False,PP_ALIGN.RIGHT)
    # 막대차트 시각화 (간단한 도형으로)
    months_data = [("10월",13,12,25),("11월",21,9,30),("12월",25,13,38)]
    max_v = 40
    bar_w = 0.55; gap = 0.15; group_w = bar_w*2+gap+0.5; start_x = 1.2
    for mi,(month,v1,v2,total) in enumerate(months_data):
        gx = start_x + mi*3.1
        bh1=2.4*(v1/max_v); bh2=2.4*(v2/max_v)
        add_rect(sl,gx,1.27+2.4-bh1,bar_w,bh1,BLUE)
        add_rect(sl,gx+bar_w+gap,1.27+2.4-bh2,bar_w,bh2,SKYLT)
        add_text(sl,str(v1),gx,1.27+2.4-bh1-0.25,bar_w,0.25,9,DARK,True,PP_ALIGN.CENTER)
        add_text(sl,str(v2),gx+bar_w+gap,1.27+2.4-bh2-0.25,bar_w,0.25,9,DARK,True,PP_ALIGN.CENTER)
        add_text(sl,month,gx-0.1,3.67+0.05,bar_w*2+gap+0.2,0.25,9.5,GRAY,False,PP_ALIGN.CENTER)
        add_text(sl,"합계 "+str(total)+"건",gx-0.1,3.92,bar_w*2+gap+0.2,0.25,9,BLUE,True,PP_ALIGN.CENTER)
    # 범례
    add_rect(sl,0.5,3.82,0.25,0.14,BLUE); add_text(sl,"1호선 EL",0.8,3.79,1.2,0.2,9,DARK)
    add_rect(sl,0.5,4.0,0.25,0.14,SKYLT); add_text(sl,"2호선 EL",0.8,3.97,1.2,0.2,9,DARK)

    months_cards=[("10월",13,12,25,"유사 수준 출발"),("11월",21,9,30,"1호선 급증"),("12월",25,13,38,"1호선 최고치")]
    for i,(mon,l1,l2,tot,comment) in enumerate(months_cards):
        x=0.2+i*3.27
        card(sl,x,4.18,3.18,1.22,AMBER if i==2 else BLUE)
        add_text(sl,mon,x+0.12,4.26,0.7,0.6,18,AMBER if i==2 else BLUE,True,PP_ALIGN.LEFT,"Calibri")
        add_text(sl,"합계 "+str(tot)+"건",x+0.8,4.26,1.0,0.3,13,DARK,True)
        add_text(sl,"1호선 "+str(l1)+"건  ·  2호선 "+str(l2)+"건",x+0.8,4.57,2.1,0.24,9,GRAY)
        add_text(sl,comment,x+0.1,4.87,2.95,0.28,9.5,"445577",False,PP_ALIGN.LEFT)

    add_rect(sl,0.2,5.43,9.6,0.18,ICE,BLUE,0.3)
    add_text(sl,"▶  12월 1호선 25건 분기 최고치  |  4분기 증가 추세 → 동절기 예방점검 강화 필요",
             0.25,5.43,9.5,0.18,8.5,NAVY,False,PP_ALIGN.LEFT)

    # ── 슬라이드 4: 장치별 분석 ───────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl, "03", "장치별 장애 분석", "EL 1호선 "+str(s1el)+"건 · 2호선 "+str(s2el)+"건")

    el1_data=[("안전스위치",14),("전장품",13),("도어",11),("제어장치",11),("제동장치",10),("속도조정",8),("부속장치",3)]
    el2_data=[("안전스위치",21),("도어",5),("속도조정장치",3),("제어장치",3),("부속장치",1),("전장품",1)]

    for side,(data,color,title,cnt) in enumerate([
        (el1_data,BLUE,"1호선 EL  장치별 분포",str(s1el)+"건"),
        (el2_data,SKY, "2호선 EL  장치별 분포",str(s2el)+"건")
    ]):
        ox = 0.2 if side==0 else 5.1; bw = 4.65 if side==0 else 4.7
        card(sl,ox,0.85,bw,4.58,color)
        add_text(sl,title,ox+0.08,0.93,bw-0.1,0.32,12,color,True)
        add_text(sl,"(총 "+cnt+")",ox+0.08,1.22,bw-0.1,0.24,9,GRAY)
        max_v=data[0][1] if data else 1
        bar_area_w=bw-1.5
        for j,(name,val) in enumerate(data):
            ry=1.5+j*0.44
            bw2=bar_area_w*(val/max_v)
            add_rect(sl,ox+1.3,ry+0.06,bar_area_w,0.28,LGRAY)
            add_rect(sl,ox+1.3,ry+0.06,max(bw2,0.05),0.28,color)
            add_text(sl,name,ox+0.08,ry,1.2,0.38,9.5,DARK,False,PP_ALIGN.RIGHT)
            add_text(sl,str(val)+"건",ox+1.35+bw2,ry+0.04,0.6,0.3,8.5,color,True)

    add_rect(sl,0.2,5.45,9.6,0.18,"FEF3C7",AMBER,0.3)
    add_text(sl,"⚠  공통 최다 : 안전스위치  —  1호선 23.7%, 2호선 61.8%  |  도어·전장품·제어장치 집중 점검 필요",
             0.25,5.45,9.5,0.18,9,"92400E",True,PP_ALIGN.LEFT)

    # ── 슬라이드 5: 역사별 현황 ───────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl, "04", "역사별 장애 현황 TOP 5")

    for side,(data,color,title) in enumerate([
        ([("월촌",11),("화원",6),("중앙로",5),("방촌",4),("영대병원",4)],BLUE,"1호선 EL  역사별 TOP 5"),
        ([("고산",5),("수성알파시티",4),("임당",3),("경대병원",3),("반월당2",3)],SKY,"2호선 EL  역사별 TOP 5"),
    ]):
        ox=0.2 if side==0 else 5.1; bw=4.65 if side==0 else 4.7
        card(sl,ox,0.85,bw,4.62,color)
        add_text(sl,title,ox+0.08,0.93,bw-0.1,0.35,12,color,True)
        max_v=data[0][1] if data else 1
        for j,(name,val) in enumerate(data):
            ry=1.38+j*0.78
            add_rect(sl,ox+0.25,ry,bw-0.3,0.7,ICE if j==0 else WHITE,LGRAY,0.3)
            badge_c=AMBER if j==0 else ("94A3B8" if j==1 else LGRAY)
            sh=sl.shapes.add_shape(1,Inches(ox+0.3),Inches(ry+0.17),Inches(0.32),Inches(0.32))
            sh.fill.solid(); sh.fill.fore_color.rgb=rgb(badge_c); sh.line.fill.background()
            add_text(sl,str(j+1),ox+0.3,ry+0.17,0.32,0.32,10,WHITE,True,PP_ALIGN.CENTER)
            add_text(sl,name,ox+0.68,ry+0.04,1.85,0.33,11,DARK,j==0)
            add_text(sl,str(val)+"건",ox+0.68,ry+0.37,0.6,0.25,9.5,color if j==0 else GRAY,True)
            bw3=3.0*(val/max_v)
            add_rect(sl,ox+1.3,ry+0.46,3.0,0.12,LGRAY)
            add_rect(sl,ox+1.3,ry+0.46,max(bw3,0.05),0.12,AMBER if j==0 else color)

    add_rect(sl,0.2,5.5,9.6,0.17,ICE,BLUE,0.3)
    add_text(sl,"▶  1호선 월촌역 11건 최다  ·  안전스위치·속도조정·도어  |  2호선 고산역 5건 최다  ·  안전스위치 집중",
             0.25,5.5,9.5,0.17,8.5,NAVY,False,PP_ALIGN.LEFT)

    # ── 슬라이드 6: 옥내/옥외 ────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl, "05", "옥내 / 옥외 장애 현황")

    for side,(outdoor,indoor,color,title) in enumerate([
        (40,19,BLUE,"1호선 EL  —  총 "+str(s1el)+"건"),
        (19,15,SKY, "2호선 EL  —  총 "+str(s2el)+"건"),
    ]):
        ox=0.2 if side==0 else 5.1; bw=4.65 if side==0 else 4.7
        total=outdoor+indoor; outpct=round(outdoor/total*100)
        card(sl,ox,0.85,bw,4.62,color)
        add_text(sl,title,ox+0.08,0.93,bw-0.1,0.32,12,color,True)
        outw=(bw-0.3)*(outdoor/total)
        add_rect(sl,ox+0.15,1.5,bw-0.3,0.45,LGRAY)
        add_rect(sl,ox+0.15,1.5,max(outw,0.05),0.45,color)
        add_text(sl,f"옥외형 {outpct}%",ox+0.15,1.97,outw,0.25,8.5,color,True,PP_ALIGN.CENTER)
        add_text(sl,f"옥내형 {100-outpct}%",ox+0.15+outw,1.97,bw-0.3-outw,0.25,8.5,GRAY,True,PP_ALIGN.CENTER)
        add_rect(sl,ox+0.15,2.35,bw/2-0.2,0.9,color)
        add_text(sl,"옥외형\n"+str(outdoor)+"건",ox+0.15,2.35,bw/2-0.2,0.9,22,WHITE,True,PP_ALIGN.CENTER)
        add_rect(sl,ox+bw/2,2.35,bw/2-0.2,0.9,"E2E8F0","CBD5E1",0.5)
        add_text(sl,"옥내형\n"+str(indoor)+"건",ox+bw/2,2.35,bw/2-0.2,0.9,22,DARK,True,PP_ALIGN.CENTER)

    add_rect(sl,0.2,5.5,9.6,0.17,ICE,BLUE,0.3)
    add_text(sl,"▶  1·2호선 모두 옥외형 비율이 높음  |  옥외 설치 EL 집중 관리 및 환경 내구성 점검 강화 필요",
             0.25,5.5,9.5,0.17,8.5,NAVY,False,PP_ALIGN.LEFT)

    # ── 슬라이드 7: 감소 대책 ────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl, "06", "장치별 원인 및 감소 대책")

    measures=[
        ("① 안전스위치",f"1호선 14건 · 2호선 21건",AMBER,
         "센서·스위치 간격 이탈, 접점 불량, 안전계통 라인 에러",
         "점검 시 적정 간격 유지·고정상태 확인, 안전라인 접점부 정밀 점검, 센서류 예비품 확보"),
        ("② 도어장치","1호선 11건 · 2호선 5건",BLUE,
         "도어 기판 에러, 스위치 간격 이탈, 롤러·세이프티슈 마모",
         "기판 노후 적기 교체, 도어 폭 측정 및 인터록 조정, 센서류 정기 점검"),
        ("③ 전장품","1호선 13건 · 2호선 1건",BLUE,
         "전자접촉기·마그네트·파워서플라이 접점 불량",
         "접점류 노후상태 점검, 파워서플라이 전압 주기 확인, 예비품 적기 교체"),
        ("④ 제어장치","1호선 11건 · 2호선 3건",BLUE,
         "제어 기판(OPB·PCB·ARD) 이상, 에러코드 발생",
         "기판 데이터 재입력 및 에러 소거, 결선상태 확인, 단자 접촉불량 점검"),
        ("⑤ 제동장치","1호선 10건 · 2호선 0건",BLUE,
         "브레이크 플런저 개방 불량, 브레이크 스위치 간격 이탈",
         "브레이크 플런저 작동 간격 조정·확인, 마그네트 접점 상태 수시 점검"),
        ("⑥ 속도조정장치","1호선 8건 · 2호선 3건",BLUE,
         "인버터 에러, 엔코더 펄스 이상, 메인보드 불량",
         "인버터 셋팅값·단자 확인, 환기팬 정상 작동 점검, 오토튜닝·층고 측정"),
    ]
    for i,(title,cnt,c,cause,action) in enumerate(measures):
        col=i%2; row=i//2
        x=0.2+col*4.95; y=0.88+row*1.55
        card(sl,x,y,4.75,1.47,c)
        add_rect(sl,x+0.07,y,4.68,0.36,c+"18",c,0.3)
        add_text(sl,title,x+0.15,y,2.6,0.36,11,c,True)
        add_text(sl,cnt,x+2.75,y,1.9,0.36,8.5,c,False,PP_ALIGN.RIGHT)
        add_text(sl,"원 인",x+0.15,y+0.38,0.55,0.26,8.5,GRAY,True)
        add_text(sl,cause,x+0.7,y+0.38,3.95,0.26,8.5,DARK)
        add_rect(sl,x+0.12,y+0.68,4.52,0.03,LGRAY)
        add_text(sl,"대 책",x+0.15,y+0.73,0.55,0.55,8.5,c,True)
        add_text(sl,action,x+0.7,y+0.73,3.95,0.65,8,DARK)

    # ── 슬라이드 8: 결론 ─────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl,0,0,10,5.625,NAVY)
    add_rect(sl,0,0,10,0.07,AMBER)
    add_rect(sl,0,0.07,0.5,5.555,BLUE)
    add_text(sl,"결론 및 중점 관리 방향",0.7,0.55,6,0.45,13,SKYLT,True)
    conclusions=[
        ("안전스위치 집중 관리",f"1호선 14건, 2호선 21건 — 양 호선 공통 최다. 정기 간격 조정·접점 점검 체계 수립",AMBER),
        ("도어·전장품 예방 교체","1호선 도어(11건)·전장품(13건) 지속 — 경년 노후 부품 교체 계획 및 예비품 확보",SKY),
        ("1호선 12월 급증 대응","10→11→12월 13→21→25건 지속 증가 — 동절기 기온 변화 예방점검 강화",SKY),
        ("월촌역 집중 점검","1호선 월촌역 11건(전체 18.6%) — 특별 점검팀 투입 검토",SKYLT),
        ("미가동 시간 단축 지속","1호선 77.42h(전년比 -9.41h) — 자재수급 지연 해소로 추가 단축 가능",GREEN),
    ]
    for i,(t,b,c) in enumerate(conclusions):
        y=1.2+i*0.76
        sh=sl.shapes.add_shape(1,Inches(0.7),Inches(y),Inches(0.38),Inches(0.38))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()
        add_text(sl,str(i+1),0.7,y,0.38,0.38,11,NAVY,True,PP_ALIGN.CENTER)
        add_text(sl,t,1.15,y,5.5,0.25,12,c,True)
        add_text(sl,b,1.15,y+0.26,5.5,0.38,9.5,"B0C8E0")
    add_rect(sl,0,5.22,10,0.405,"070F2B")
    add_text(sl,"2025년 4분기  엘리베이터 장애분석 보고서  |  대구도시철도공사 시설운영처",
             0.5,5.24,9,0.35,10,"4A7099",False,PP_ALIGN.CENTER)

    out=io.BytesIO(); prs.save(out); out.seek(0)
    return out.read()
