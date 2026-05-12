"""
ES 장애분석 보고서 생성 (python-pptx 기반)
Vercel Serverless 환경 - 외부 프로세스 없음
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

def rgb(h):
    h=h.lstrip('#')
    return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

NAVY="0B2535"; TEAL="0D7966"; TEALLT="14B8A6"; TEALX="5EEAD4"
ICE="CCFBF1"; WHITE="FFFFFF"; OFFWH="F0FDF9"; GRAY="64748B"
LGRAY="E2E8F0"; DARK="1E293B"; AMBER="F59E0B"; GREEN="059669"
BLUE="1D4ED8"; RED="DC2626"; REDLT="FEF2F2"

W=Inches(10); H=Inches(5.625)

def add_rect(sl,x,y,w,h,fc,lc=None,lw=None):
    sh=sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fc)
    if lc: sh.line.color.rgb=rgb(lc)
    if lw: sh.line.width=Pt(lw)
    else:
        if not lc: sh.line.fill.background()
    return sh

def add_text(sl,text,x,y,w,h,size=12,color=WHITE,bold=False,
             align=PP_ALIGN.LEFT,font="Malgun Gothic"):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(text)
    r.font.size=Pt(size); r.font.color.rgb=rgb(color)
    r.font.bold=bold; r.font.name=font
    return tb

def add_header(sl,num,title,sub=None):
    add_rect(sl,0,0,10,0.68,NAVY)
    add_rect(sl,0,0,0.07,0.68,AMBER)
    add_rect(sl,0,0.68,10,0.045,TEAL)
    add_text(sl,num,0.15,0.05,0.55,0.58,10,AMBER,True,PP_ALIGN.CENTER,"Calibri")
    add_text(sl,title,0.75,0.08,8.5,0.55,19,WHITE,True)
    if sub: add_text(sl,sub,6.5,0.08,3.4,0.55,10,TEALX,False,PP_ALIGN.RIGHT)

def card(sl,x,y,w,h,ac=None):
    add_rect(sl,x,y,w,h,WHITE,LGRAY,0.5)
    if ac: add_rect(sl,x,y,w,0.06,ac)

def build_es_report(stats):
    prs=Presentation()
    prs.slide_width=W; prs.slide_height=H

    s1es=stats.get('1호선ES',0)
    s2es=stats.get('2호선ES',0)
    total_es=s1es+s2es

    # ── 슬라이드 1: 표지 ─────────────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl,0,0,10,5.625,NAVY)
    add_rect(sl,0,0,10,0.07,AMBER)
    add_rect(sl,0,0.07,0.5,5.555,TEAL)
    for ox,oy,ow,oh,oc in [(6.0,-0.4,4.8,4.8,TEAL),(6.9,0.5,3.0,3.0,TEALLT),(7.5,1.1,1.8,1.8,TEALX)]:
        sh=sl.shapes.add_shape(9,Inches(ox),Inches(oy),Inches(ow),Inches(oh))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(oc)
        sh.line.color.rgb=rgb(WHITE); sh.line.width=Pt(0.5)
    add_rect(sl,0.7,1.2,1.95,0.38,TEAL)
    add_text(sl,"ESCALATOR  ·  ES",0.7,1.22,1.95,0.34,9,WHITE,True,PP_ALIGN.CENTER,"Calibri")
    add_rect(sl,0.7,1.65,2.4,0.34,AMBER)
    add_text(sl,"2025년  4분기",0.7,1.67,2.4,0.3,12,NAVY,True,PP_ALIGN.CENTER)
    add_text(sl,"에스컬레이터",0.7,2.12,9,0.95,46,WHITE,True,PP_ALIGN.LEFT,"Calibri")
    add_text(sl,"장애분석 보고서",0.7,3.0,9,0.9,38,TEALX,True,PP_ALIGN.LEFT,"Calibri")
    add_text(sl,"1호선 · 2호선  |  대구도시철도공사 시설운영처",0.7,4.0,8.5,0.38,12,"7ECFC4")
    add_rect(sl,0,5.22,10,0.405,"040E18")
    add_text(sl,f"ES 1호선 {s1es}건 · ES 2호선 {s2es}건  |  합계 {total_es}건",
             0.5,5.24,9,0.35,10.5,"3D9E90",False,PP_ALIGN.CENTER)

    # ── 슬라이드 2: 종합 현황 ─────────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl,"01","종합 현황","2025년 4분기 (10월 ~ 12월)")
    kpis=[
        ("1호선 ES 장애건수",str(s1es)+"건","전분기 130건 대비 -8건",TEAL,GREEN),
        ("2호선 ES 장애건수",str(s2es)+"건","전분기 147건 대비 +31건",TEAL,RED),
        ("평균 복구시간","88분","1호선 108분 · 2호선 68분",GRAY,GRAY),
        ("총 미가동시간","304h","1호선 181h · 2호선 123h",AMBER,GREEN),
    ]
    for i,(label,val,sub,c,tc) in enumerate(kpis):
        x=0.2+i*2.45
        card(sl,x,0.85,2.35,1.52,c)
        add_text(sl,label,x,0.93,2.35,0.28,9.5,c,True,PP_ALIGN.CENTER)
        add_text(sl,val,x,1.18,2.35,0.65,30,DARK,True,PP_ALIGN.CENTER,"Calibri")
        add_text(sl,sub,x,1.83,2.35,0.28,8.5,tc,False,PP_ALIGN.CENTER)

    # 1호선 상세
    card(sl,0.2,2.5,4.65,2.82,TEAL)
    add_text(sl,"1호선 ES  —  "+str(s1es)+"건",0.28,2.58,4.5,0.32,12,TEAL,True)
    for j,(n,v,p) in enumerate([("안전스위치","56건","45.9%"),("제어장치","14건","11.5%"),("핸드레일","13건","10.7%")]):
        ry=2.93+j*0.66
        add_rect(sl,2.5,ry,2.25,0.58,ICE if j==0 else WHITE,LGRAY,0.3)
        sh=sl.shapes.add_shape(9,Inches(2.55),Inches(ry+0.1),Inches(0.35),Inches(0.35))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(TEAL if j==0 else LGRAY); sh.line.fill.background()
        add_text(sl,str(j+1),2.55,ry+0.1,0.35,0.35,10,WHITE,True,PP_ALIGN.CENTER)
        add_text(sl,n,2.95,ry+0.02,1.75,0.28,10,DARK,j==0)
        add_text(sl,v+"("+p+")",2.95,ry+0.3,1.75,0.22,8.5,GRAY)

    # 2호선 상세
    card(sl,5.1,2.5,4.7,2.82,TEALLT)
    add_text(sl,"2호선 ES  —  "+str(s2es)+"건",5.18,2.58,4.55,0.32,12,TEAL,True)
    for j,(n,v,p) in enumerate([("안전스위치","80건","44.9%"),("부속장치","34건","19.1%"),("동력전달장치","20건","11.2%")]):
        ry=2.93+j*0.66
        add_rect(sl,7.38,ry,2.28,0.58,ICE if j==0 else WHITE,LGRAY,0.3)
        sh=sl.shapes.add_shape(9,Inches(7.43),Inches(ry+0.1),Inches(0.35),Inches(0.35))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(TEALLT if j==0 else LGRAY); sh.line.fill.background()
        add_text(sl,str(j+1),7.43,ry+0.1,0.35,0.35,10,WHITE,True,PP_ALIGN.CENTER)
        add_text(sl,n,7.83,ry+0.02,1.78,0.28,10,DARK,j==0)
        add_text(sl,v+"("+p+")",7.83,ry+0.3,1.78,0.22,8.5,GRAY)

    add_rect(sl,0.2,5.35,9.6,0.22,"FEF3C7",AMBER,0.5)
    add_text(sl,"⚠  1·2호선 공통 : 안전스위치 최다 (1호선 45.9% · 2호선 44.9%)  |  2호선 부속장치 19.1% 특이 높음",
             0.25,5.35,9.5,0.22,9,"92400E",True,PP_ALIGN.LEFT)

    # ── 슬라이드 3: 월별 현황 ─────────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl,"02","월별 장애 현황","10월 · 11월 · 12월")
    card(sl,0.2,0.85,9.6,3.2,TEAL)
    add_text(sl,"ES 1·2호선 월별 장애건수 추이",0.28,0.93,6,0.32,12,TEAL,True)
    add_text(sl,"1호선: 37→40→45건  ·  2호선: 68→49→61건",5.5,0.93,4.25,0.32,9,GRAY,False,PP_ALIGN.RIGHT)
    months_data=[("10월",37,68,105),("11월",40,49,89),("12월",45,61,106)]
    max_v=80
    for mi,(month,v1,v2,total) in enumerate(months_data):
        gx=1.2+mi*3.1
        bh1=2.3*(v1/max_v); bh2=2.3*(v2/max_v)
        add_rect(sl,gx,1.27+2.3-bh1,0.55,bh1,TEAL)
        add_rect(sl,gx+0.7,1.27+2.3-bh2,0.55,bh2,TEALLT)
        add_text(sl,str(v1),gx,1.27+2.3-bh1-0.25,0.55,0.25,9,DARK,True,PP_ALIGN.CENTER)
        add_text(sl,str(v2),gx+0.7,1.27+2.3-bh2-0.25,0.55,0.25,9,DARK,True,PP_ALIGN.CENTER)
        add_text(sl,month,gx-0.1,3.57+0.05,1.3,0.25,9.5,GRAY,False,PP_ALIGN.CENTER)
        add_text(sl,"합계 "+str(total)+"건",gx-0.1,3.82,1.3,0.25,9,TEAL,True,PP_ALIGN.CENTER)
    add_rect(sl,0.5,3.82,0.25,0.14,TEAL); add_text(sl,"1호선 ES",0.8,3.79,1.2,0.2,9,DARK)
    add_rect(sl,0.5,4.0,0.25,0.14,TEALLT); add_text(sl,"2호선 ES",0.8,3.97,1.2,0.2,9,DARK)

    months_cards=[("10월",37,68,105,"2호선 집중 발생"),("11월",40,49,89,"분기 최저"),("12월",45,61,106,"1호선 최고치")]
    for i,(mon,l1,l2,tot,comment) in enumerate(months_cards):
        x=0.2+i*3.27
        card(sl,x,4.18,3.18,1.22,TEALLT if i==0 else TEAL)
        add_text(sl,mon,x+0.12,4.26,0.7,0.6,18,TEALLT if i==0 else TEAL,True,PP_ALIGN.LEFT,"Calibri")
        add_text(sl,"합계 "+str(tot)+"건",x+0.8,4.26,1.0,0.3,13,DARK,True)
        add_text(sl,"1호선 "+str(l1)+"건  ·  2호선 "+str(l2)+"건",x+0.8,4.57,2.1,0.24,9,GRAY)
        add_text(sl,comment,x+0.1,4.87,2.95,0.28,9.5,TEAL)

    add_rect(sl,0.2,5.43,9.6,0.18,ICE,TEAL,0.3)
    add_text(sl,"▶  2호선 10월 68건 최다  |  1호선 꾸준한 증가세 (37→45건)  |  전분기 277건 대비 +23건",
             0.25,5.43,9.5,0.18,8.5,NAVY,False,PP_ALIGN.LEFT)

    # ── 슬라이드 4: 장치별 분석 ───────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl,"03","장치별 장애 분석","ES 1호선 "+str(s1es)+"건 · 2호선 "+str(s2es)+"건")

    es1_data=[("안전스위치",56),("제어장치",14),("핸드레일",13),("동력전달장치",9),("제동장치",9),("속도조정장치",9),("스텝콤",7),("전장품",3),("부속장치",2)]
    es2_data=[("안전스위치",80),("부속장치",34),("동력전달장치",20),("스텝콤",14),("제동장치",12),("속도조정장치",8),("핸드레일",4),("동력장치",3),("제어장치",2)]

    for side,(data,color,title,cnt) in enumerate([
        (es1_data,TEAL,"1호선 ES  장치별 분포",str(s1es)+"건"),
        (es2_data,TEALLT,"2호선 ES  장치별 분포",str(s2es)+"건")
    ]):
        ox=0.2 if side==0 else 5.1; bw=4.65 if side==0 else 4.7
        card(sl,ox,0.85,bw,4.58,color)
        add_text(sl,title,ox+0.08,0.93,bw-0.1,0.32,12,color,True)
        add_text(sl,"(총 "+cnt+")",ox+0.08,1.22,bw-0.1,0.24,9,GRAY)
        max_v=data[0][1] if data else 1
        bar_area_w=bw-1.5
        for j,(name,val) in enumerate(data):
            ry=1.5+j*0.37
            bw2=bar_area_w*(val/max_v)
            add_rect(sl,ox+1.3,ry+0.05,bar_area_w,0.24,LGRAY)
            add_rect(sl,ox+1.3,ry+0.05,max(bw2,0.05),0.24,color)
            add_text(sl,name,ox+0.08,ry,1.2,0.33,9,DARK,False,PP_ALIGN.RIGHT)
            add_text(sl,str(val)+"건",ox+1.35+bw2,ry+0.02,0.6,0.25,8.5,color,True)

    add_rect(sl,0.2,5.45,9.6,0.18,"FEF3C7",AMBER,0.3)
    add_text(sl,"⚠  공통 최다 : 안전스위치 (1호선 45.9% · 2호선 44.9%)  |  2호선 부속장치(34건)·동력전달(20건) 집중 점검",
             0.25,5.45,9.5,0.18,9,"92400E",True,PP_ALIGN.LEFT)

    # ── 슬라이드 5: 역사별 현황 ───────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl,"04","역사별 장애 현황 TOP 5")

    for side,(data,color,title) in enumerate([
        ([("설화명곡",17),("화원",15),("각산",10),("진천",9),("신기",9)],TEAL,"1호선 ES  역사별 TOP 5"),
        ([("죽전",14),("임당",14),("이곡",13),("청라언덕2",10),("담티",10)],TEALLT,"2호선 ES  역사별 TOP 5"),
    ]):
        ox=0.2 if side==0 else 5.1; bw=4.65 if side==0 else 4.7
        card(sl,ox,0.85,bw,4.62,color)
        add_text(sl,title,ox+0.08,0.93,bw-0.1,0.35,12,color,True)
        max_v=data[0][1] if data else 1
        for j,(name,val) in enumerate(data):
            ry=1.38+j*0.78
            add_rect(sl,ox+0.25,ry,bw-0.3,0.7,ICE if j==0 else WHITE,LGRAY,0.3)
            badge_c=AMBER if j==0 else ("94A3B8" if j==1 else LGRAY)
            sh=sl.shapes.add_shape(1,Inches(ox+0.3),Inches(ry+0.17),Inches(0.32),Inches(0.32))
            sh.fill.solid(); sh.fill.fore_color.rgb=rgb(badge_c); sh.line.fill.background()
            add_text(sl,str(j+1),ox+0.3,ry+0.17,0.32,0.32,10,WHITE,True,PP_ALIGN.CENTER)
            add_text(sl,name,ox+0.68,ry+0.04,1.85,0.33,11,DARK,j==0)
            add_text(sl,str(val)+"건",ox+0.68,ry+0.37,0.6,0.25,9.5,color if j==0 else GRAY,True)
            bw3=3.0*(val/max_v)
            add_rect(sl,ox+1.3,ry+0.46,3.0,0.12,LGRAY)
            add_rect(sl,ox+1.3,ry+0.46,max(bw3,0.05),0.12,AMBER if j==0 else color)

    add_rect(sl,0.2,5.5,9.6,0.17,ICE,TEAL,0.3)
    add_text(sl,"▶  1호선 설화명곡(17건)·화원(15건) 집중  |  2호선 죽전·임당 공동 1위(14건)  — 상위 3개역 중점 관리",
             0.25,5.5,9.5,0.17,8.5,NAVY,False,PP_ALIGN.LEFT)

    # ── 슬라이드 6: 옥내/옥외 ────────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl,"05","옥내 / 옥외 장애 현황")

    for side,(outdoor,indoor,color,title,special) in enumerate([
        (73,49,TEAL,"1호선 ES  —  총 "+str(s1es)+"건",False),
        (83,95,TEALLT,"2호선 ES  —  총 "+str(s2es)+"건",True),
    ]):
        ox=0.2 if side==0 else 5.1; bw=4.65 if side==0 else 4.7
        total=outdoor+indoor; outpct=round(outdoor/total*100)
        card(sl,ox,0.85,bw,4.62,color)
        add_text(sl,title,ox+0.08,0.93,bw-0.1,0.32,12,color,True)
        if special:
            add_text(sl,"★ 옥내형이 옥외형 초과",ox+0.08,1.22,bw-0.1,0.22,9.5,RED,True)
        outw=(bw-0.3)*(outdoor/total)
        add_rect(sl,ox+0.15,1.55,bw-0.3,0.42,LGRAY)
        add_rect(sl,ox+0.15,1.55,max(outw,0.05),0.42,color)
        add_text(sl,f"옥외형 {outpct}%",ox+0.15,2.0,outw,0.25,8.5,color,True,PP_ALIGN.CENTER)
        add_text(sl,f"옥내형 {100-outpct}%",ox+0.15+outw,2.0,bw-0.3-outw,0.25,8.5,GRAY,True,PP_ALIGN.CENTER)
        add_rect(sl,ox+0.15,2.38,bw/2-0.2,0.88,color)
        add_text(sl,"옥외형\n"+str(outdoor)+"건",ox+0.15,2.38,bw/2-0.2,0.88,22,WHITE,True,PP_ALIGN.CENTER)
        add_rect(sl,ox+bw/2,2.38,bw/2-0.2,0.88,"E2E8F0","CBD5E1",0.5)
        add_text(sl,"옥내형\n"+str(indoor)+"건",ox+bw/2,2.38,bw/2-0.2,0.88,22,DARK,True,PP_ALIGN.CENTER)

    add_rect(sl,0.2,5.5,9.6,0.17,"FEF3C7",AMBER,0.5)
    add_text(sl,"⚠  2호선 ES : 옥내형(95건)이 옥외형(83건) 초과 — 지하역사 내부 에스컬레이터 집중 점검 필요",
             0.25,5.5,9.5,0.17,9,"92400E",True,PP_ALIGN.LEFT)

    # ── 슬라이드 7: 감소 대책 ────────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_header(sl,"06","장치별 원인 및 감소 대책")

    measures=[
        ("① 안전스위치","1호선 56건 · 2호선 80건",AMBER,
         "안전스위치·센서류 간격 이탈, 접점불량, 안전계통 라인 에러",
         "점검 시 스위치·센서류 적정 간격 유지, 고정상태 확인, 접점부 정밀 점검"),
        ("② 부속장치","1호선 2건 · 2호선 34건",TEAL,
         "플레이트 이상, 스텝레일 파손, 스커트·콤플레이트 장애",
         "이상 발견 즉시 신속 보수, 정기 교체 계획 수립"),
        ("③ 동력전달장치","1호선 9건 · 2호선 20건",TEAL,
         "스텝체인·베어링 장력 이완, 구동축·스프라켓 이상",
         "구동체인 장력 조정, 주유상태 유지, 압력벨트·구동휠 마모 확인"),
        ("④ 제동장치","1호선 9건 · 2호선 12건",TEAL,
         "보조브레이크 에러, 전자 브레이크 작동불량, 라이닝 마모",
         "브레이크 간격·센서 적절 조정, 전압 상태 점검"),
        ("⑤ 핸드레일","1호선 13건 · 2호선 4건",TEAL,
         "핸드레일 구동부품 불량, 장력 이완, 곡각부 베어링 마모",
         "장력 적정 조정, 구동드라이브·곡각부 베어링 적기 교체"),
        ("⑥ 스텝·콤","1호선 7건 · 2호선 14건",TEAL,
         "스텝롤러 파손, 콤·데마케이션 이상",
         "파손 스텝·콤·데마케이션 적기 교체, 스텝롤러 정기 점검"),
    ]
    for i,(title,cnt,c,cause,action) in enumerate(measures):
        col=i%2; row=i//2
        x=0.2+col*4.95; y=0.88+row*1.55
        card(sl,x,y,4.75,1.47,c)
        add_rect(sl,x+0.07,y,4.68,0.36,c+"18",c,0.3)
        add_text(sl,title,x+0.15,y,2.6,0.36,11,c,True)
        add_text(sl,cnt,x+2.75,y,1.9,0.36,8.5,c,False,PP_ALIGN.RIGHT)
        add_text(sl,"원 인",x+0.15,y+0.38,0.55,0.26,8.5,GRAY,True)
        add_text(sl,cause,x+0.7,y+0.38,3.95,0.26,8.5,DARK)
        add_rect(sl,x+0.12,y+0.68,4.52,0.03,LGRAY)
        add_text(sl,"대 책",x+0.15,y+0.73,0.55,0.55,8.5,c,True)
        add_text(sl,action,x+0.7,y+0.73,3.95,0.65,8,DARK)

    # ── 슬라이드 8: 결론 ─────────────────────────────────
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl,0,0,10,5.625,NAVY)
    add_rect(sl,0,0,10,0.07,AMBER)
    add_rect(sl,0,0.07,0.5,5.555,TEAL)
    add_text(sl,"결론 및 중점 관리 방향",0.7,0.55,6,0.45,13,TEALX,True)
    conclusions=[
        ("안전스위치 정밀 관리 체계 구축",f"양 호선 공통 최다(1호선 56건·2호선 80건) — 전체 45% 차지. 정기 점검 체계 수립 필수",AMBER),
        ("2호선 부속장치 34건 — 특별 관리","전체의 19.1%로 이례적 높은 비율 — 스텝레일·플레이트 노후화 원인 분석 및 계획 교체",TEALLT),
        ("설화명곡·화원역 집중 점검","1호선 설화명곡 17건·화원 15건 — 전체의 26.2% 집중. 중점 관리역 지정 및 집중 순회점검",TEALLT),
        ("2호선 옥내형 초과 — 실내 집중","옥내형 95건 > 옥외형 83건(유일한 역전) — 지하역사 내부 에스컬레이터 환경 점검 강화",TEALLT),
        ("미가동 시간 감소 추세 유지","1호선 180.87h(전분기比 -221.83h) — 개선 성과 지속 관리 필요","059669"),
    ]
    for i,(t,b,c) in enumerate(conclusions):
        y=1.2+i*0.76
        sh=sl.shapes.add_shape(1,Inches(0.7),Inches(y),Inches(0.38),Inches(0.38))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()
        add_text(sl,str(i+1),0.7,y,0.38,0.38,11,NAVY,True,PP_ALIGN.CENTER)
        add_text(sl,t,1.15,y,5.5,0.25,12,c,True)
        add_text(sl,b,1.15,y+0.26,5.5,0.38,9.5,"7ECFC4")
    add_rect(sl,0,5.22,10,0.405,"040E18")
    add_text(sl,"2025년 4분기  에스컬레이터 장애분석 보고서  |  대구도시철도공사 시설운영처",
             0.5,5.24,9,0.35,10,"2D7A6D",False,PP_ALIGN.CENTER)

    out=io.BytesIO(); prs.save(out); out.seek(0)
    return out.read()
